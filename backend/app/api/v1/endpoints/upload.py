import os
import time
import json

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, BackgroundTasks, Request
from sqlalchemy.orm import Session

from app.config import settings
from app.core.dependencies import get_admin_user
from app.models.database import get_db
from app.models.user import User
from app.models.presentation import Presentation
from app.models.upload_log import UploadLog
from app.schemas.presentation import PresentationResponse
from app.services.extraction.pipeline import extract_presentation
from app.services.extraction.progress import init_progress, complete_progress, fail_progress, cleanup_progress
from app.services.html_generator import generate_webpage

router = APIRouter(prefix="/admin", tags=["admin"])


def process_pptx(presentation_id: int, file_path: str, media_dir: str, slide_data_path: str):
    """Background task to extract PPTX content."""
    from pptx import Presentation as PptxPresentation
    from app.models.database import SessionLocal

    db = SessionLocal()
    start_time = time.time()
    try:
        # Count total slides and initialize progress tracking
        prs = PptxPresentation(file_path)
        total_slides = len(prs.slides)
        init_progress(presentation_id, total_slides)

        result = extract_presentation(file_path, presentation_id, media_dir)
        with open(slide_data_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False)

        # Generate HTML webpage using Claude API (best-effort)
        pres_dir = os.path.dirname(slide_data_path)
        try:
            update_progress = __import__(
                "app.services.extraction.progress", fromlist=["update_progress"]
            ).update_progress
            update_progress(
                presentation_id,
                current_slide=len(result["slides"]),
                phase="webpage",
                message="Generating web page with AI...",
            )
            webpage_path = generate_webpage(
                presentation_id, slide_data_path, media_dir, pres_dir
            )
            if webpage_path:
                import logging
                logging.getLogger(__name__).info(
                    "Webpage generated: %s", webpage_path
                )
        except Exception as e:
            import logging
            logging.getLogger(__name__).warning(
                "Webpage generation failed (non-fatal): %s", e
            )

        presentation = db.query(Presentation).filter(Presentation.id == presentation_id).first()
        presentation.status = "ready"
        presentation.slide_count = len(result["slides"])
        presentation.slide_width_emu = result["slide_width_emu"]
        presentation.slide_height_emu = result["slide_height_emu"]
        presentation.title = result.get("title") or presentation.title

        log = db.query(UploadLog).filter(UploadLog.presentation_id == presentation_id).first()
        if log:
            log.status = "success"
            log.processing_time_ms = int((time.time() - start_time) * 1000)

        db.commit()
        complete_progress(presentation_id, message=f"Extraction complete! {total_slides} slides processed.")
    except Exception as e:
        fail_progress(presentation_id, error=str(e))

        presentation = db.query(Presentation).filter(Presentation.id == presentation_id).first()
        if presentation:
            presentation.status = "failed"
            presentation.error_message = str(e)

        log = db.query(UploadLog).filter(UploadLog.presentation_id == presentation_id).first()
        if log:
            log.status = "failed"
            log.error_message = str(e)
            log.processing_time_ms = int((time.time() - start_time) * 1000)

        db.commit()
    finally:
        db.close()


@router.post("/upload", response_model=PresentationResponse, status_code=202)
async def upload_presentation(
    background_tasks: BackgroundTasks,
    request: Request,
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    user: User = Depends(get_admin_user),
):
    if not file.filename.lower().endswith((".pptx",)):
        raise HTTPException(status_code=400, detail="Only .pptx files are allowed")

    content = await file.read()
    file_size = len(content)

    if file_size > settings.MAX_UPLOAD_SIZE_MB * 1024 * 1024:
        raise HTTPException(status_code=400, detail=f"File too large. Max size: {settings.MAX_UPLOAD_SIZE_MB}MB")

    title = os.path.splitext(file.filename)[0]

    presentation = Presentation(
        title=title,
        original_filename=file.filename,
        file_path="",
        slide_data_path="",
        media_dir="",
        slide_count=0,
        status="processing",
        uploaded_by=user.id,
    )
    db.add(presentation)
    db.flush()

    pres_dir = os.path.join(settings.STORAGE_DIR, "presentations", str(presentation.id))
    media_dir = os.path.join(pres_dir, "media")
    upload_dir = os.path.join(settings.STORAGE_DIR, "uploads", str(presentation.id))
    os.makedirs(media_dir, exist_ok=True)
    os.makedirs(upload_dir, exist_ok=True)

    file_path = os.path.join(upload_dir, "original.pptx")
    with open(file_path, "wb") as f:
        f.write(content)

    slide_data_path = os.path.join(pres_dir, "slides.json")

    presentation.file_path = file_path
    presentation.slide_data_path = slide_data_path
    presentation.media_dir = media_dir

    log = UploadLog(
        presentation_id=presentation.id,
        original_filename=file.filename,
        file_size_bytes=file_size,
        status="processing",
        uploaded_by=user.id,
        ip_address=request.client.host if request.client else None,
        user_agent=request.headers.get("user-agent"),
    )
    db.add(log)
    db.commit()
    db.refresh(presentation)

    background_tasks.add_task(process_pptx, presentation.id, file_path, media_dir, slide_data_path)

    return presentation
