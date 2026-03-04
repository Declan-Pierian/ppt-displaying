"""Main extraction pipeline: converts a PPTX file into structured JSON."""

import os
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from app.services.extraction.text_extractor import extract_text_body
from app.services.extraction.image_extractor import extract_image_shape, extract_background
from app.services.extraction.table_extractor import extract_table_shape
from app.services.extraction.chart_extractor import extract_chart_shape
from app.services.extraction.shape_extractor import extract_auto_shape
from app.services.extraction.group_extractor import extract_group_shape
from app.services.extraction.utils import get_fill_data, get_border_data
from app.services.extraction.progress import update_progress
from app.services.extraction.image_exporter import export_slide_images

logger = logging.getLogger(__name__)


def extract_presentation(pptx_path: str, presentation_id: int, media_dir: str) -> dict:
    """Main entry point: extract all content from a PPTX file."""
    prs = Presentation(pptx_path)
    title = _extract_title(prs)

    result = {
        "presentation_id": presentation_id,
        "title": title,
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": [],
    }

    for slide_index, slide in enumerate(prs.slides):
        logger.info(f"Extracting slide {slide_index + 1}/{len(prs.slides)}")

        slide_data = {
            "slide_index": slide_index,
            "slide_number": slide_index + 1,
            "background": extract_background(slide, media_dir, slide_index, presentation_id),
            "shapes": [],
            "notes": _extract_notes(slide),
        }

        for z_order, shape in enumerate(slide.shapes):
            try:
                shape_data = extract_shape(shape, media_dir, z_order, presentation_id)
                if shape_data:
                    slide_data["shapes"].append(shape_data)
            except Exception as e:
                logger.warning(f"Failed to extract shape on slide {slide_index + 1}: {e}")
                try:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                        base = _make_base(shape, z_order)
                        shape_data = _extract_text_shape(shape, base)
                        if shape_data:
                            slide_data["shapes"].append(shape_data)
                except Exception:
                    pass

        result["slides"].append(slide_data)
        update_progress(
            presentation_id,
            current_slide=slide_index + 1,
            phase="extraction",
            message=f"Extracted slide {slide_index + 1} of {len(prs.slides)}",
        )

    # Guarantee correct slide ordering
    result["slides"].sort(key=lambda s: s["slide_index"])

    # --- Render each slide as a high-res PNG image (graceful degradation) ---
    try:
        slide_images_dir = os.path.join(media_dir, "slide_images")
        os.makedirs(slide_images_dir, exist_ok=True)

        exported_paths = export_slide_images(pptx_path, slide_images_dir)

        # Map each exported image back to its slide data
        for img_path in exported_paths:
            filename = os.path.basename(img_path)  # e.g. "slide_1.png"
            # Extract slide number from filename pattern "slide_<N>.png"
            try:
                slide_num = int(filename.replace("slide_", "").replace(".png", ""))
                slide_idx = slide_num - 1
                for slide_data in result["slides"]:
                    if slide_data["slide_index"] == slide_idx:
                        slide_data["slide_image"] = f"media/slide_images/{filename}"
                        break
            except (ValueError, IndexError):
                logger.warning("Could not map exported image %s to a slide.", filename)

        logger.info("Slide image export completed: %d images generated.", len(exported_paths))
    except Exception as e:
        logger.warning("Slide image rendering failed (extraction data is still intact): %s", e)

    return result


def _make_base(shape, z_order: int) -> dict:
    """Build the base shape data dict."""
    return {
        "shape_id": f"sp_{shape.shape_id}",
        "position": {
            "left_emu": shape.left if shape.left is not None else 0,
            "top_emu": shape.top if shape.top is not None else 0,
            "width_emu": shape.width if shape.width is not None else 0,
            "height_emu": shape.height if shape.height is not None else 0,
        },
        "rotation_degrees": shape.rotation if shape.rotation else 0,
        "z_order": z_order,
    }


def _has_image(shape) -> bool:
    """Safely check if a shape contains an image."""
    try:
        # Method 1: Direct image access
        if hasattr(shape, 'image'):
            try:
                _ = shape.image.content_type
                return True
            except Exception:
                pass

        # Method 2: XML inspection for blipFill/pic elements
        elem = shape._element
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

        tag = elem.tag
        if tag.endswith('}pic') or 'pic' in tag.lower():
            return True

        # Look for blipFill with an actual embedded reference
        blip_fill = elem.find(f'.//{ns_a}blipFill')
        if blip_fill is not None:
            blip = blip_fill.find(f'{ns_a}blip')
            if blip is not None and blip.get(f'{ns_r}embed'):
                return True

        # Broader search through all descendants
        for child in elem.iter():
            if child.tag.endswith('}blip'):
                embed = child.get(f'{ns_r}embed')
                if embed:
                    return True

    except Exception:
        pass
    return False


def _has_chart(shape) -> bool:
    """Safely check if a shape contains a chart."""
    try:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            return True
        if hasattr(shape, 'chart'):
            _ = shape.chart.chart_type
            return True
    except Exception:
        pass
    return False


def _has_table(shape) -> bool:
    """Safely check if a shape contains a table."""
    try:
        if hasattr(shape, 'has_table') and shape.has_table:
            return True
    except Exception:
        pass
    return False


def extract_shape(shape, media_dir: str, z_order: int, presentation_id: int) -> dict | None:
    """Dispatch shape extraction based on shape type."""
    base = _make_base(shape, z_order)

    try:
        shape_type = shape.shape_type
    except Exception:
        shape_type = None

    # --- GROUP ---
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return extract_group_shape(shape, base, media_dir, presentation_id, extract_shape)

    # --- PICTURE ---
    if shape_type == MSO_SHAPE_TYPE.PICTURE or (shape_type is None and _has_image(shape)):
        try:
            return extract_image_shape(shape, base, media_dir, presentation_id)
        except Exception as e:
            logger.warning(f"Image extraction failed: {e}")

    # --- TABLE ---
    if shape_type == MSO_SHAPE_TYPE.TABLE or _has_table(shape):
        try:
            return extract_table_shape(shape, base)
        except Exception as e:
            logger.warning(f"Table extraction failed: {e}")

    # --- CHART ---
    if shape_type == MSO_SHAPE_TYPE.CHART or _has_chart(shape):
        try:
            return extract_chart_shape(shape, base, media_dir)
        except Exception as e:
            logger.warning(f"Chart extraction failed: {e}")

    # --- MEDIA ---
    if shape_type == MSO_SHAPE_TYPE.MEDIA:
        return _extract_media_shape(shape, base, media_dir, presentation_id)

    # --- PLACEHOLDER ---
    if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return _extract_placeholder_content(shape, base, media_dir, presentation_id)

    # --- AUTO_SHAPE ---
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return extract_auto_shape(shape, base)

    # --- FREEFORM ---
    if shape_type == MSO_SHAPE_TYPE.FREEFORM:
        return extract_auto_shape(shape, base)

    # --- TEXT BOX ---
    if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return _extract_text_shape(shape, base)

    # --- EMBEDDED OLE ---
    if shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        if _has_image(shape):
            try:
                return extract_image_shape(shape, base, media_dir, presentation_id)
            except Exception:
                pass
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
            return _extract_text_shape(shape, base)
        return None

    # --- SMART FALLBACK ---
    if _has_image(shape):
        try:
            return extract_image_shape(shape, base, media_dir, presentation_id)
        except Exception:
            pass

    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            return _extract_text_shape(shape, base)

    # Try auto shape only if it has visible content
    try:
        result = extract_auto_shape(shape, base)
        if result:
            fill = result.get("fill", {})
            text_body = result.get("text_body")
            has_text = text_body and text_body.get("paragraphs") and any(
                p.get("runs") for p in text_body["paragraphs"]
            )
            has_fill = fill and fill.get("type") not in (None, "none")
            if has_text or has_fill:
                return result
    except Exception:
        pass

    return None


def _extract_placeholder_content(shape, base: dict, media_dir: str, presentation_id: int) -> dict | None:
    """Smart extraction for placeholder shapes."""
    if _has_image(shape):
        try:
            return extract_image_shape(shape, base, media_dir, presentation_id)
        except Exception:
            pass

    if _has_chart(shape):
        try:
            return extract_chart_shape(shape, base, media_dir)
        except Exception:
            pass

    if _has_table(shape):
        try:
            return extract_table_shape(shape, base)
        except Exception:
            pass

    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            return _extract_text_shape(shape, base)

    return None


def _extract_text_shape(shape, base: dict) -> dict:
    """Extract a text box or text-containing placeholder."""
    fill_data = {"type": "none"}
    try:
        fill_data = get_fill_data(shape.fill)
    except Exception:
        pass

    border_data = {"color": None, "width_pt": None, "dash_style": None}
    try:
        border_data = get_border_data(shape.line)
    except Exception:
        pass

    text_body = extract_text_body(shape.text_frame)
    hyperlink = _extract_shape_hyperlink(shape)

    base.update({
        "shape_type": "text_box",
        "fill": fill_data,
        "border": border_data,
        "text_body": text_body,
        "hyperlink": hyperlink,
    })
    return base


def _extract_shape_hyperlink(shape) -> str | None:
    """Extract hyperlink from a shape's click action."""
    try:
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        click_action = shape._element.find(f'.//{ns_a}hlinkClick')
        if click_action is not None:
            r_id = click_action.get(f'{ns_r}id')
            if r_id:
                rel = shape.part.rels.get(r_id)
                if rel and hasattr(rel, 'target_ref'):
                    return rel.target_ref
    except Exception:
        pass
    return None


def _extract_media_shape(shape, base: dict, media_dir: str, presentation_id: int) -> dict:
    """Extract a media (video/audio) shape."""
    import os

    media_path = None
    try:
        sp_elem = shape._element
        ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        for child in sp_elem.iter():
            r_link = child.get(f'{ns_r}link') or child.get(f'{ns_r}embed')
            if r_link:
                try:
                    rel = shape.part.rels.get(r_link)
                    if rel and hasattr(rel, 'target_part'):
                        blob = rel.target_part.blob
                        ext = os.path.splitext(rel.target_ref)[1] if hasattr(rel, 'target_ref') else '.mp4'
                        filename = f"media_{shape.shape_id}{ext}"
                        filepath = os.path.join(media_dir, filename)
                        with open(filepath, "wb") as f:
                            f.write(blob)
                        media_path = f"media/{filename}"
                        break
                except Exception:
                    continue
    except Exception:
        pass

    base.update({
        "shape_type": "media",
        "media": {"media_path": media_path, "media_type": "video"},
    })
    return base


def _extract_title(prs) -> str:
    """Extract presentation title from the first slide."""
    try:
        if prs.slides and len(prs.slides) > 0:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                        try:
                            if shape.placeholder_format.idx in (0, 15):
                                text = shape.text_frame.text.strip()
                                if text:
                                    return text
                        except Exception:
                            pass
            for shape in first_slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.text.strip()[:100]
    except Exception:
        pass
    return "Untitled Presentation"


def _extract_notes(slide) -> str | None:
    """Extract speaker notes from a slide."""
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            return notes_slide.notes_text_frame.text.strip() or None
    except Exception:
        pass
    return None
