"""Extract text content with full formatting from text frames."""

from pptx.util import Pt
from app.services.extraction.utils import get_color_hex, emu_to_pt


def extract_text_body(text_frame) -> dict:
    """Extract full text body with paragraphs, runs, and formatting."""
    if text_frame is None:
        return {"paragraphs": []}

    paragraphs = []
    for para in text_frame.paragraphs:
        p_data = _extract_paragraph(para)
        paragraphs.append(p_data)

    return {"paragraphs": paragraphs}


def _extract_paragraph(para) -> dict:
    """Extract a single paragraph with all its properties."""
    alignment = "left"
    try:
        if para.alignment is not None:
            alignment = str(para.alignment).split(".")[-1].split("(")[0].strip().lower()
    except Exception:
        pass

    line_spacing_pt = None
    try:
        if para.line_spacing is not None:
            if isinstance(para.line_spacing, (int, float)):
                line_spacing_pt = para.line_spacing
            else:
                line_spacing_pt = emu_to_pt(para.line_spacing)
    except Exception:
        pass

    space_before_pt = None
    try:
        if para.space_before is not None:
            space_before_pt = emu_to_pt(para.space_before)
    except Exception:
        pass

    space_after_pt = None
    try:
        if para.space_after is not None:
            space_after_pt = emu_to_pt(para.space_after)
    except Exception:
        pass

    level = 0
    try:
        level = para.level if para.level is not None else 0
    except Exception:
        pass

    bullet = _extract_bullet(para)
    runs = []
    for run in para.runs:
        r_data = _extract_run(run)
        runs.append(r_data)

    # If no runs but paragraph has text, create a synthetic run
    if not runs and para.text.strip():
        runs.append({
            "text": para.text,
            "font": {
                "name": None,
                "size_pt": None,
                "bold": False,
                "italic": False,
                "underline": False,
                "strikethrough": False,
                "color": None,
            },
            "hyperlink": None,
        })

    return {
        "alignment": alignment,
        "line_spacing_pt": line_spacing_pt,
        "space_before_pt": space_before_pt,
        "space_after_pt": space_after_pt,
        "level": level,
        "bullet": bullet,
        "runs": runs,
    }


def _extract_run(run) -> dict:
    """Extract a single text run with font properties."""
    font = run.font

    font_name = None
    try:
        font_name = font.name
    except Exception:
        pass

    size_pt = None
    try:
        if font.size is not None:
            size_pt = round(font.size.pt, 1)
    except Exception:
        pass

    bold = False
    try:
        bold = bool(font.bold) if font.bold is not None else False
    except Exception:
        pass

    italic = False
    try:
        italic = bool(font.italic) if font.italic is not None else False
    except Exception:
        pass

    underline = False
    try:
        underline = bool(font.underline) if font.underline is not None else False
    except Exception:
        pass

    strikethrough = False
    try:
        # python-pptx doesn't directly expose strikethrough on the font object,
        # check the XML element directly
        from lxml import etree
        rPr = run._r.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
        if rPr is not None:
            strike = rPr.get('strike')
            if strike and strike != 'noStrike':
                strikethrough = True
    except Exception:
        pass

    color = get_color_hex(font.color) if font.color else None

    hyperlink_url = None
    try:
        if run.hyperlink and run.hyperlink.address:
            hyperlink_url = run.hyperlink.address
    except Exception:
        pass

    return {
        "text": run.text,
        "font": {
            "name": font_name,
            "size_pt": size_pt,
            "bold": bold,
            "italic": italic,
            "underline": underline,
            "strikethrough": strikethrough,
            "color": color,
        },
        "hyperlink": {"url": hyperlink_url} if hyperlink_url else None,
    }


def _extract_bullet(para) -> dict | None:
    """Extract bullet information from a paragraph."""
    try:
        pPr = para._p.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        if pPr is None:
            return None

        # Check for bullet character
        buChar = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        if buChar is not None:
            char = buChar.get('char', '\u2022')
            return {"type": "char", "char": char}

        # Check for bullet auto-number
        buAutoNum = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
        if buAutoNum is not None:
            num_type = buAutoNum.get('type', 'arabicPeriod')
            return {"type": "number", "number_type": num_type}

        # Check for no bullet explicitly set
        buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
        if buNone is not None:
            return None

    except Exception:
        pass

    return None
